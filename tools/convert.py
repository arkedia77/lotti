#!/usr/bin/env python3
"""PDF/PPTX → JSON 변환 도구
Usage:
  python convert.py input.pdf          → input.json
  python convert.py input.pptx         → input.json
  python convert.py input.pdf -o out.json
"""
import sys, json, os, argparse

def pdf_to_dict(path):
    import pdfplumber
    result = {"source": os.path.basename(path), "type": "pdf", "pages": []}
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            p = {"page": i + 1, "text": page.extract_text() or ""}
            tables = page.extract_tables()
            if tables:
                p["tables"] = tables
            result["pages"].append(p)
    result["total_pages"] = len(result["pages"])
    return result

def pptx_to_dict(path):
    from pptx import Presentation
    result = {"source": os.path.basename(path), "type": "pptx", "slides": []}
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        s = {"slide": i + 1, "texts": []}
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    s["texts"].append(text)
            if shape.has_table:
                table = []
                for row in shape.table.rows:
                    table.append([cell.text.strip() for cell in row.cells])
                s.setdefault("tables", []).append(table)
        result["slides"].append(s)
    result["total_slides"] = len(result["slides"])
    return result

def main():
    parser = argparse.ArgumentParser(description="PDF/PPTX → JSON")
    parser.add_argument("input", help="입력 파일 (pdf/pptx)")
    parser.add_argument("-o", "--output", help="출력 JSON 경로 (기본: 같은 이름.json)")
    args = parser.parse_args()

    ext = os.path.splitext(args.input)[1].lower()
    if ext == ".pdf":
        data = pdf_to_dict(args.input)
    elif ext in (".pptx", ".ppt"):
        data = pptx_to_dict(args.input)
    else:
        print(f"지원하지 않는 형식: {ext}", file=sys.stderr)
        sys.exit(1)

    out = args.output or os.path.splitext(args.input)[0] + ".json"
    with open(out, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"변환 완료: {out} ({data.get('total_pages') or data.get('total_slides')} 페이지/슬라이드)")

if __name__ == "__main__":
    main()
